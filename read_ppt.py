import sys
from pptx import Presentation

def main():
    ppt_path = r"C:\Users\kinga\Downloads\Presentation1.pptx"
    try:
        prs = Presentation(ppt_path)
        print(f"Presentation has {len(prs.slides)} slides.")
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"  - Text: {shape.text[:50]}")
                if shape.shape_type == 13: # Picture
                    print(f"  - Picture found")
    except Exception as e:
        print(f"Error reading presentation: {e}")

if __name__ == "__main__":
    main()
