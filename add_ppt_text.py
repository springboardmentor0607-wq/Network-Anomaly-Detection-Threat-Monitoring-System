import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def parse_markdown(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    slides_data = []
    # Split by slide headers
    slide_sections = re.split(r'## Slide \d+:', content)
    
    for section in slide_sections[1:]: # Skip the first part before any Slide header
        lines = section.strip().split('\n')
        title = lines[0].strip()
        bullets = []
        for line in lines[1:]:
            line = line.strip()
            if not line or line.startswith('**(') or line.startswith('*Tip:'):
                continue
            # Remove Markdown formatting
            clean_line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            clean_line = re.sub(r'_(.*?)_', r'\1', clean_line)
            clean_line = re.sub(r'`(.*?)`', r'\1', clean_line)
            bullets.append(clean_line)
        slides_data.append({'title': title, 'bullets': bullets})
    return slides_data

def add_text_to_ppt(ppt_path, out_path, md_path):
    slides_data = parse_markdown(md_path)
    prs = Presentation(ppt_path)
    
    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break
        slide = prs.slides[i]
        
        # Add a title text box
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data['title']
        p.font.bold = True
        p.font.size = Pt(28)
        
        # Add content text box
        contentBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        tf_content = contentBox.text_frame
        tf_content.word_wrap = True
        
        for j, bullet in enumerate(slide_data['bullets']):
            if j == 0:
                p_bullet = tf_content.paragraphs[0]
            else:
                p_bullet = tf_content.add_paragraph()
            
            p_bullet.text = bullet
            if bullet.startswith('* '):
                p_bullet.level = 0
                p_bullet.text = bullet[2:]
            elif bullet.startswith('*'):
                p_bullet.level = 1
                p_bullet.text = bullet[1:].strip()
            elif bullet.startswith('- '):
                p_bullet.level = 1
                p_bullet.text = bullet[2:]
            else:
                p_bullet.level = 0
                
            p_bullet.font.size = Pt(14)
            
    prs.save(out_path)
    print(f"Successfully added text to the slides. Saved to {out_path}")

if __name__ == "__main__":
    ppt = r"C:\Users\kinga\Downloads\Presentation1.pptx"
    out = r"C:\Users\kinga\Downloads\Presentation1_with_text.pptx"
    md = r"e:\NetShield\Presentation_Outline.md"
    add_text_to_ppt(ppt, out, md)
